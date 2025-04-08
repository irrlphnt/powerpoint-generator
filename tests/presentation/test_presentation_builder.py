import pytest
from pathlib import Path
from pptx import Presentation
from src.presentation import PresentationBuilder

class TestPresentationBuilder:
    @pytest.fixture
    def builder(self):
        """Create a fresh PresentationBuilder instance for each test."""
        return PresentationBuilder()
    
    def test_initialization(self, builder):
        """Test that the PresentationBuilder initializes correctly."""
        assert builder.presentation is not None
        assert builder.current_slide is None
        assert builder.template is None
        assert builder.theme is None
        
    def test_set_template(self, builder):
        """Test setting a template."""
        template_name = "business"
        builder.set_template(template_name)
        assert builder.template == template_name
        
    def test_set_theme(self, builder):
        """Test setting a theme."""
        theme_name = "modern"
        builder.set_theme(theme_name)
        assert builder.theme == theme_name
        
    def test_add_title_slide(self, builder):
        """Test adding a title slide."""
        title = "Test Title"
        subtitle = "Test Subtitle"
        
        builder.add_title_slide(title, subtitle)
        assert builder.current_slide is not None
        
        # Verify slide content
        assert builder.current_slide.shapes.title.text == title
        assert builder.current_slide.placeholders[1].text == subtitle
        
    def test_add_title_slide_no_subtitle(self, builder):
        """Test adding a title slide without subtitle."""
        title = "Test Title"
        
        builder.add_title_slide(title)
        assert builder.current_slide is not None
        assert builder.current_slide.shapes.title.text == title
        
    def test_add_content_slide(self, builder):
        """Test adding a content slide with bullet points."""
        title = "Test Content"
        content = ["Point 1", "Point 2", "Point 3"]
        
        builder.add_content_slide(title, content)
        assert builder.current_slide is not None
        
        # Verify slide content
        assert builder.current_slide.shapes.title.text == title
        text_frame = builder.current_slide.placeholders[1].text_frame
        assert len(text_frame.paragraphs) == len(content)
        for i, paragraph in enumerate(text_frame.paragraphs):
            assert paragraph.text == content[i]
            
    def test_add_image_slide(self, builder, tmp_path):
        """Test adding an image slide."""
        # Create a temporary image file for testing
        image_path = tmp_path / "test_image.jpg"
        image_path.write_bytes(b"fake image data")
        
        title = "Test Image"
        caption = "Test Caption"
        
        builder.add_image_slide(title, str(image_path), caption)
        assert builder.current_slide is not None
        
        # Verify slide content
        assert builder.current_slide.shapes.title.text == title
        # Note: We can't easily verify the image was added correctly without
        # more complex testing of the pptx library internals
        
    def test_save_presentation(self, builder, tmp_path):
        """Test saving a presentation to a file."""
        # Add some content
        builder.add_title_slide("Test Title", "Test Subtitle")
        
        # Save to temporary file
        output_path = tmp_path / "test_presentation.pptx"
        builder.save(str(output_path))
        
        # Verify file was created
        assert output_path.exists()
        
        # Verify we can open it with python-pptx
        presentation = Presentation(str(output_path))
        assert len(presentation.slides) == 1
        
    def test_multiple_slides(self, builder):
        """Test creating a presentation with multiple slides."""
        # Add multiple slides
        builder.add_title_slide("Title", "Subtitle")
        builder.add_content_slide("Content", ["Point 1", "Point 2"])
        
        # Verify slide count
        assert len(builder.presentation.slides) == 2
        
    def test_invalid_image_path(self, builder):
        """Test handling of invalid image path."""
        with pytest.raises(Exception):
            builder.add_image_slide("Test", "nonexistent.jpg")
            
    def test_empty_content_slide(self, builder):
        """Test creating a content slide with empty content."""
        builder.add_content_slide("Empty", [])
        assert builder.current_slide is not None
        assert builder.current_slide.shapes.title.text == "Empty"
        assert len(builder.current_slide.placeholders[1].text_frame.paragraphs) == 0 
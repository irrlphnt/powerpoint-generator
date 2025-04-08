from typing import List, Optional, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class PresentationBuilder:
    """Main class for building PowerPoint presentations."""
    
    def __init__(self):
        """Initialize a new presentation builder."""
        self.presentation = Presentation()
        self.current_slide = None
        self.template = None
        self.theme = None
        
    def set_template(self, template_name: str) -> None:
        """Set the template for the presentation.
        
        Args:
            template_name: Name of the template to use
        """
        self.template = template_name
        # TODO: Implement template loading
        
    def set_theme(self, theme_name: str) -> None:
        """Set the theme for the presentation.
        
        Args:
            theme_name: Name of the theme to use
        """
        self.theme = theme_name
        # TODO: Implement theme application
        
    def add_title_slide(self, title: str, subtitle: Optional[str] = None) -> None:
        """Add a title slide to the presentation.
        
        Args:
            title: Main title text
            subtitle: Optional subtitle text
        """
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        self.current_slide = self.presentation.slides.add_slide(slide_layout)
        
        title_shape = self.current_slide.shapes.title
        title_shape.text = title
        
        if subtitle:
            subtitle_shape = self.current_slide.placeholders[1]
            subtitle_shape.text = subtitle
            
    def add_content_slide(self, title: str, content: List[str], layout: str = "bullet_list") -> None:
        """Add a content slide with bullet points.
        
        Args:
            title: Slide title
            content: List of bullet points
            layout: Layout type (default: "bullet_list")
        """
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        self.current_slide = self.presentation.slides.add_slide(slide_layout)
        
        title_shape = self.current_slide.shapes.title
        title_shape.text = title
        
        content_shape = self.current_slide.placeholders[1]
        text_frame = content_shape.text_frame
        
        for point in content:
            paragraph = text_frame.add_paragraph()
            paragraph.text = point
            paragraph.level = 0
            
    def add_image_slide(self, title: str, image_path: str, caption: Optional[str] = None) -> None:
        """Add a slide with an image.
        
        Args:
            title: Slide title
            image_path: Path to the image file
            caption: Optional caption text
        """
        slide_layout = self.presentation.slide_layouts[5]  # Title only layout
        self.current_slide = self.presentation.slides.add_slide(slide_layout)
        
        title_shape = self.current_slide.shapes.title
        title_shape.text = title
        
        # Add image
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        self.current_slide.shapes.add_picture(image_path, left, top, width=width)
        
        if caption:
            left = Inches(1)
            top = Inches(6)
            width = Inches(8)
            height = Inches(0.5)
            textbox = self.current_slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            paragraph = text_frame.add_paragraph()
            paragraph.text = caption
            paragraph.alignment = PP_ALIGN.CENTER
            
    def save(self, output_path: str) -> None:
        """Save the presentation to a file.
        
        Args:
            output_path: Path where to save the presentation
        """
        self.presentation.save(output_path) 